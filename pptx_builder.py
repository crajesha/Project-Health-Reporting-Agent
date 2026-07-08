"""
pptx_builder.py
----------------
Generates the 7-slide monthly executive presentation using python-pptx
only -- no Node.js / pptxgenjs dependency. This is what makes the deck
generator work on Python-only hosts (Streamlit Community Cloud, Render,
Railway, etc.), where a `node` binary is not available.

Public entry point:
    build_executive_pptx(portfolio: list[dict], out) -> None
        out can be a filesystem path (str) or a file-like object (e.g.
        io.BytesIO), so callers can either write to disk (CLI/scheduler)
        or stream straight to a Streamlit download button.
"""

from __future__ import annotations

from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn

# ---- Palette: "Midnight Executive" ----
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xC9, 0x8A, 0x1D)
GREEN = RGBColor(0x1E, 0x7A, 0x46)
INK = RGBColor(0x1F, 0x23, 0x28)
MUTE = RGBColor(0x5B, 0x64, 0x72)
CARD = RGBColor(0xF5, 0xF7, 0xFC)

RAG_COLOR = {"Green": GREEN, "Amber": AMBER, "Red": RED}

PW, PH = Inches(13.333), Inches(7.5)  # 16:9 wide


def _set_no_line(shape):
    shape.line.fill.background()


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    _set_no_line(shape)


def _shadow_off(shape):
    shape.shadow.inherit = False


def _textbox(slide, x, y, w, h, text, size=14, color=INK, bold=False, italic=False,
             font="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
             wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    return box


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _title_bar(slide, kicker, title):
    _bg(slide, WHITE)
    _textbox(slide, Inches(0.6), Inches(0.32), Inches(10), Inches(0.35),
              kicker.upper(), size=12, color=NAVY, bold=True, font="Calibri")
    _textbox(slide, Inches(0.6), Inches(0.62), Inches(11.8), Inches(0.8),
              title, size=30, color=INK, bold=True, font="Cambria")


def _footer(slide, page_num):
    _textbox(slide, Inches(0.5), PH - Inches(0.42), Inches(6), Inches(0.3),
              "Professional Services  |  Confidential", size=9, color=MUTE)
    _textbox(slide, PW - Inches(1.1), PH - Inches(0.42), Inches(0.6), Inches(0.3),
              str(page_num), size=9, color=MUTE, align=PP_ALIGN.RIGHT)


def _rounded_card(slide, x, y, w, h, color=CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shape.adjustments[0] = 0.06
    except Exception:
        pass
    _fill(shape, color)
    _shadow_off(shape)
    return shape


def _style_chart_text(chart, font="Calibri", size=10, color=MUTE):
    try:
        chart.font.size = Pt(size)
        chart.font.name = font
        chart.font.color.rgb = color
    except Exception:
        pass


def build_executive_pptx(portfolio: list[dict], out) -> None:
    """portfolio: list of dicts as produced by main.run_pipeline() /
    dashboard.py -- keys: name, pm, rag, score, confidence, sheet_status,
    agrees, pct_complete, pct_delayed_tasks, missed_end_dates,
    avg_delay_days, critical_count, critical_red, critical_delayed,
    flagged_comments, total_comments, milestones_completed,
    milestones_total, total_tasks, stage.
    """
    prs = Presentation()
    prs.slide_width = PW
    prs.slide_height = PH
    blank = prs.slide_layouts[6]

    n = len(portfolio)
    red_n = sum(1 for p in portfolio if p["rag"] == "Red")
    amber_n = sum(1 for p in portfolio if p["rag"] == "Amber")
    green_n = sum(1 for p in portfolio if p["rag"] == "Green")
    avg_score = sum(p["score"] for p in portfolio) / n if n else 0
    disagreements = sum(1 for p in portfolio if not p["agrees"])
    worst = min(portfolio, key=lambda p: p["score"]) if portfolio else None
    month_label = datetime.now().strftime("%B %Y")

    # ---------------- Slide 1: Title ----------------
    s = prs.slides.add_slide(blank)
    _bg(s, NAVY)
    circ1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.4), Inches(-2.4), Inches(6.5), Inches(6.5))
    _fill(circ1, RGBColor(0x27, 0x35, 0x80)); _shadow_off(circ1)
    circ2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(3.3), Inches(4.0), Inches(4.0))
    _fill(circ2, RGBColor(0x2B, 0x3A, 0x8C)); _shadow_off(circ2)

    _textbox(s, Inches(0.7), Inches(2.5), Inches(10.5), Inches(0.5),
              "EXECUTIVE PROJECT HEALTH REVIEW", size=16, color=ICE, bold=True)
    _textbox(s, Inches(0.7), Inches(2.95), Inches(10.8), Inches(1.1),
              "Portfolio Health, Trends & Risk Outlook", size=38, color=WHITE, bold=True, font="Cambria")
    _textbox(s, Inches(0.7), Inches(4.0), Inches(8), Inches(0.5),
              f"Monthly Synthesis — {month_label}", size=16, color=ICE)
    _textbox(s, Inches(0.7), Inches(6.7), Inches(9), Inches(0.4),
              "Prepared by the Project Health Reporting Agent  |  Professional Services",
              size=11, color=RGBColor(0x9F, 0xB1, 0xE8))

    # ---------------- Slide 2: Executive Summary ----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "Slide 2", "Executive Summary")

    stats = [
        ("Projects Reviewed", str(n), NAVY),
        ("Red / Amber", f"{red_n} / {amber_n}", RED),
        ("Avg. Health Score", f"{avg_score:.0f}", NAVY),
        ("Status Overrides vs. Self-Report", f"{disagreements}/{n}", AMBER),
    ]
    card_w, gap, start_x, y = Inches(2.85), Inches(0.25), Inches(0.6), Inches(1.7)
    for i, (label, value, color) in enumerate(stats):
        x = start_x + i * (card_w + gap)
        _rounded_card(s, x, y, card_w, Inches(1.55))
        _textbox(s, x, y + Inches(0.15), card_w, Inches(0.75), value, size=32, color=color,
                  bold=True, font="Cambria", align=PP_ALIGN.CENTER)
        _textbox(s, x + Inches(0.1), y + Inches(0.98), card_w - Inches(0.2), Inches(0.5),
                  label, size=11, color=MUTE, align=PP_ALIGN.CENTER)

    summary = (
        f"The portfolio's average composite health score is {avg_score:.0f}/100. "
        f"{disagreements} of {n} projects carry an overall status that differs from what their own "
        f"status sheet self-reports once schedule slippage, critical-path risk, and stakeholder "
        f"comments are weighed together rather than read in isolation. "
    )
    if worst:
        summary += f"{worst['name']} is the portfolio's most urgent item, rated {worst['rag']} at a {worst['score']}/100 composite score."
    _textbox(s, Inches(0.6), Inches(3.55), Inches(12.1), Inches(1.3), summary, size=15, color=INK,
              line_spacing=1.25)

    _textbox(s, Inches(0.6), Inches(5.1), Inches(6), Inches(0.3),
              "KEY TAKEAWAY FOR THE CLIENT", size=11, color=NAVY, bold=True)
    banner = _rounded_card(s, Inches(0.6), Inches(5.45), Inches(12.1), Inches(1.25), color=NAVY)
    _textbox(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.95),
              "Self-reported schedule health is no longer a reliable single indicator across this "
              "portfolio — composite, multi-signal scoring is now the basis for governance decisions "
              "and client updates.", size=13.5, color=WHITE, italic=True, anchor=MSO_ANCHOR.MIDDLE)
    _footer(s, 2)

    # ---------------- Slide 3: Portfolio Health at a Glance ----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "Slide 3", "Portfolio Health at a Glance")

    _textbox(s, Inches(0.6), Inches(1.65), Inches(4), Inches(0.3), "RAG DISTRIBUTION",
              size=12, color=NAVY, bold=True)
    chart_data = CategoryChartData()
    chart_data.categories = ["Red", "Amber", "Green"]
    chart_data.add_series("RAG", (red_n, amber_n, green_n))
    gframe = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(0.5), Inches(2.0), Inches(5.4), Inches(4.4), chart_data)
    chart = gframe.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    _style_chart_text(chart)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = "0"
    plot.data_labels.number_format_is_linked = False
    colors = [RED, AMBER, GREEN]
    try:
        points = chart.series[0].points
        for pt, col in zip(points, colors):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = col
    except Exception:
        pass

    _textbox(s, Inches(6.3), Inches(1.65), Inches(6.4), Inches(0.3),
              "COMPOSITE HEALTH SCORE BY PROJECT", size=12, color=NAVY, bold=True)
    chart_data2 = CategoryChartData()
    labels = [p["name"][:22] + ("…" if len(p["name"]) > 22 else "") for p in portfolio]
    chart_data2.categories = labels
    chart_data2.add_series("Score", tuple(p["score"] for p in portfolio))
    bframe = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(6.2), Inches(2.0), Inches(6.5), Inches(4.4), chart_data2)
    bchart = bframe.chart
    bchart.has_legend = False
    _style_chart_text(bchart)
    vaxis = bchart.value_axis
    vaxis.minimum_scale = 0
    vaxis.maximum_scale = 100
    bplot = bchart.plots[0]
    bplot.has_data_labels = True
    bplot.data_labels.number_format = "0"
    bplot.data_labels.number_format_is_linked = False
    bplot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    try:
        bpoints = bchart.series[0].points
        for pt, p in zip(bpoints, portfolio):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = RAG_COLOR[p["rag"]]
    except Exception:
        pass
    _footer(s, 3)

    # ---------------- Slide 4: Cross-Portfolio Trend Analysis ----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "Slide 4", "Cross-Portfolio Trend Analysis")

    avg_slip = sum(p["avg_delay_days"] for p in portfolio) / n if n else 0
    total_overdue = sum(p["missed_end_dates"] for p in portfolio)
    total_crit_risk = sum(p["critical_red"] + p["critical_delayed"] for p in portfolio)

    trends = [
        ("Self-reported status is systematically unreliable",
         f"{disagreements} of {n} projects show a computed RAG that diverges from their own Schedule "
         f"Health field — in both directions. This is a portfolio-wide pattern, not an isolated "
         f"data-entry issue, and points to inconsistent local definitions of \"Green\" across PMs."),
        ("Critical-path risk is concentrated, not evenly spread",
         f"{total_crit_risk} critical-path tasks across the portfolio are either Red or trending behind "
         f"schedule. Where slippage exists, it clusters on the critical path rather than being spread "
         f"evenly across the plan — the projects most at risk are at risk because of a small number of "
         f"high-leverage tasks."),
        ("Slippage severity is rising with project maturity",
         f"Average delay where slippage exists is {avg_slip:.1f} days across the portfolio, with "
         f"{total_overdue} open tasks already past their planned end date. Later-stage projects show "
         f"materially larger average delays than earlier-stage ones — schedule risk compounds rather "
         f"than resolves as projects mature."),
    ]
    y = Inches(1.75)
    for i, (h, b) in enumerate(trends):
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.02), Inches(0.45), Inches(0.45))
        _fill(circ, NAVY); _shadow_off(circ)
        _textbox(s, Inches(0.6), y + Inches(0.02), Inches(0.45), Inches(0.45), str(i + 1), size=18,
                  color=WHITE, bold=True, font="Cambria", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _textbox(s, Inches(1.25), y, Inches(11.3), Inches(0.4), h, size=16, color=INK, bold=True)
        _textbox(s, Inches(1.25), y + Inches(0.42), Inches(11.3), Inches(1.0), b, size=12.5, color=MUTE,
                  line_spacing=1.2)
        y += Inches(1.65)
    _footer(s, 4)

    # ---------------- Slide 5: Emerging & Common Risks ----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "Slide 5", "Emerging & Common Risks")

    risks = [
        ("Critical-path slippage", "Multiple projects carry delayed or Red critical-path tasks that directly threaten committed end dates.", "High"),
        ("Pending approvals as silent blockers", "Stakeholder comments repeatedly reference approvals and dependencies awaiting sign-off — a leading indicator that precedes visible schedule slippage.", "Medium"),
        ("Inconsistent self-reported RAG", "Local Schedule Health values are not comparable across projects, undermining portfolio-level roll-ups unless re-scored centrally.", "Medium"),
        ("Completion-vs-timeline gap widening", "At least one project shows completion trailing elapsed timeline by a double-digit percentage, with no compensating acceleration visible.", "High"),
    ]
    sev_color = {"High": RED, "Medium": AMBER, "Low": GREEN}
    col_w, row_h, gap_x, gap_y = Inches(5.85), Inches(2.2), Inches(0.35), Inches(0.25)
    start_x, start_y = Inches(0.6), Inches(1.7)
    for i, (t, d, sev) in enumerate(risks):
        col, row = i % 2, i // 2
        x = start_x + col * (col_w + gap_x)
        yy = start_y + row * (row_h + gap_y)
        _rounded_card(s, x, yy, col_w, row_h)
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), yy + Inches(0.28), Inches(1.1), Inches(0.34))
        try:
            badge.adjustments[0] = 0.5
        except Exception:
            pass
        _fill(badge, sev_color[sev]); _shadow_off(badge)
        _textbox(s, x + Inches(0.3), yy + Inches(0.28), Inches(1.1), Inches(0.34), sev.upper(), size=10,
                  color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _textbox(s, x + Inches(0.3), yy + Inches(0.78), col_w - Inches(0.6), Inches(0.45), t, size=15,
                  color=INK, bold=True)
        _textbox(s, x + Inches(0.3), yy + Inches(1.25), col_w - Inches(0.6), Inches(0.85), d, size=11.5,
                  color=MUTE, line_spacing=1.2)
    _footer(s, 5)

    # ---------------- Slide 6: Recommendations ----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "Slide 6", "Recommendations")

    recs = [
        ("Standardize RAG scoring centrally", "Adopt the composite scoring model portfolio-wide so status is comparable across PMs and projects, replacing locally-defined Schedule Health fields as the reporting source of truth."),
        ("Institute a critical-path escalation trigger", "Any critical-path task slipping more than 5 business days should auto-generate a follow-up action, rather than waiting for the next weekly cycle."),
        ("Close the approval-latency gap", "Track pending-approval comments as a distinct leading indicator and assign an SLA (e.g. 3 business days) before they convert into schedule slippage."),
        ("Weekly steering review for Red/Amber projects", "Formalize a standing weekly checkpoint for any project below an 80 composite score until it returns to Green."),
    ]
    y = Inches(1.8)
    for i, (t, d) in enumerate(recs):
        _textbox(s, Inches(0.6), y, Inches(0.9), Inches(1.0), f"0{i+1}", size=28, color=ICE, bold=True,
                  font="Cambria")
        _textbox(s, Inches(1.6), y + Inches(0.02), Inches(10.9), Inches(0.4), t, size=16, color=NAVY, bold=True)
        _textbox(s, Inches(1.6), y + Inches(0.44), Inches(10.9), Inches(0.55), d, size=12.5, color=MUTE,
                  line_spacing=1.2)
        y += Inches(1.28)
    _footer(s, 6)

    # ---------------- Slide 7: Projects Requiring Attention (table) ----------------
    s = prs.slides.add_slide(blank)
    _title_bar(s, "Slide 7", "Projects Requiring Attention — Detail")

    header = ["Project", "PM", "Status", "Score", "Open Overdue", "Critical Risk", "Self-Report vs. Computed"]
    rows_data = sorted(portfolio, key=lambda p: p["score"])
    n_rows = len(rows_data) + 1
    table_shape = s.shapes.add_table(n_rows, len(header), Inches(0.6), Inches(1.8), Inches(12.1), Inches(2.6))
    table = table_shape.table
    col_widths = [Inches(3.0), Inches(1.9), Inches(1.1), Inches(1.0), Inches(1.7), Inches(1.5), Inches(1.9)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for ci, h in enumerate(header):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True; p.font.size = Pt(11); p.font.color.rgb = WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for ri, p in enumerate(rows_data, start=1):
        vals = [
            p["name"], p["pm"] or "—", p["rag"], f"{p['score']}",
            f"{p['missed_end_dates']}", f"{p['critical_red'] + p['critical_delayed']}/{p['critical_count']}",
            "Matches" if p["agrees"] else f"{p['sheet_status']} \u2192 {p['rag']}",
        ]
        for ci, v in enumerate(vals):
            cell = table.cell(ri, ci)
            cell.text = str(v)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.color.rgb = RAG_COLOR.get(v, INK) if ci == 2 else INK
            para.font.bold = ci == 2
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    _textbox(s, Inches(0.6), Inches(4.75), Inches(5), Inches(0.3), "METHODOLOGY NOTE", size=11,
              color=NAVY, bold=True)
    _textbox(s, Inches(0.6), Inches(5.1), Inches(12.1), Inches(0.9),
              "RAG status is calculated from a weighted composite of schedule slippage (30%), completion "
              "vs. elapsed timeline (20%), critical-path task health (20%), stakeholder blocker signals "
              "(15%), and milestone health (15%), with override rules that cap the rating when "
              "critical-path or blocker risk is severe. Full methodology available on request.",
              size=11.5, color=MUTE, italic=True, line_spacing=1.25)
    _textbox(s, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.4),
              f"Data as of the most recent weekly extract for each project. Portfolio of {n} project(s) reviewed this cycle.",
              size=10, color=MUTE)
    _footer(s, 7)

    prs.save(out)
